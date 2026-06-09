from __future__ import annotations

import io
import json
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence

from src.class_registry import ClassRegistry

_STEP_TYPES = {
    "click",
    "drag",
    "validate_pins",
    "click_sequence",
    "type_text",
    "press_key",
    "wait_ms",
    "auth_sequence",
    "input_text",
    "mold_setup",
}
_DOC_SUFFIXES = {".txt", ".md", ".pdf", ".pptx"}


@dataclass
class SOPSourceRef:
    kind: str
    index: int
    label: str
    text: str

    def to_json(self) -> Dict[str, Any]:
        return {
            "kind": self.kind,
            "index": self.index,
            "label": self.label,
            "text": self.text,
        }


@dataclass
class SOPAtom:
    id: str
    name: str
    type: str
    target: Optional[str] = None
    description: str = ""
    enabled: bool = True
    class_name: Optional[str] = None
    confidence: Optional[float] = None
    source_page: Optional[int] = None
    source_span: Optional[str] = None
    preconditions: List[str] = field(default_factory=list)
    postconditions: List[str] = field(default_factory=list)
    raw_text: str = ""
    extra: Dict[str, Any] = field(default_factory=dict)

    def to_step(self) -> Dict[str, Any]:
        step: Dict[str, Any] = {
            "id": self.id,
            "name": self.name,
            "type": self.type,
            "description": self.description,
            "enabled": self.enabled,
        }
        if self.target:
            step["target"] = self.target
        if self.class_name:
            step["class_name"] = self.class_name
        if self.confidence is not None:
            step["confidence"] = round(float(self.confidence), 4)
        if self.source_page is not None:
            step["source_page"] = int(self.source_page)
        if self.source_span:
            step["source_span"] = self.source_span
        if self.preconditions:
            step["preconditions"] = list(self.preconditions)
        if self.postconditions:
            step["postconditions"] = list(self.postconditions)
        if self.raw_text:
            step["source_text"] = self.raw_text
        step.update(self.extra)
        return step


@dataclass
class SOPDocumentArtifact:
    version: str
    title: str
    source_path: str
    source_type: str
    raw_text: str
    steps: List[Dict[str, Any]]
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_json(self) -> Dict[str, Any]:
        return {
            "version": self.version,
            "title": self.title,
            "source_path": self.source_path,
            "source_type": self.source_type,
            "raw_text": self.raw_text,
            "steps": self.steps,
            "metadata": self.metadata,
        }


@dataclass
class SOPDocumentExtraction:
    source_path: str
    source_type: str
    title: str
    raw_text: str
    refs: List[SOPSourceRef] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_json(self) -> Dict[str, Any]:
        return {
            "source_path": self.source_path,
            "source_type": self.source_type,
            "title": self.title,
            "raw_text": self.raw_text,
            "refs": [ref.to_json() for ref in self.refs],
            "metadata": dict(self.metadata),
        }


class SOPDocumentIngestor:
    """Read PDF/TXT SOP files and atomize them into schema-ready JSON."""

    def __init__(self, llm: Any | None = None, ocr: Any | None = None) -> None:
        self._llm = llm
        self._ocr = ocr

    def ingest(self, source_path: str | Path) -> SOPDocumentArtifact:
        path = Path(source_path)
        if not path.exists():
            raise FileNotFoundError(f"SOP document not found: {path}")
        if path.suffix.lower() not in _DOC_SUFFIXES:
            raise ValueError(f"Unsupported SOP document type: {path.suffix}")

        extraction = self.extract_document(path)
        source_type = extraction.source_type
        raw_text = extraction.raw_text
        registry = ClassRegistry.load()
        class_names = registry.class_names()

        if self._llm is not None and raw_text.strip():
            from src.sop_llm_atomizer import SOPLLMAtomizer  # noqa: PLC0415

            atomizer = SOPLLMAtomizer(llm=self._llm)
            atomized = atomizer.atomize(extraction)
            if (
                atomized.atomization_mode in {"llm", "rules_fallback"}
                and atomized.steps
            ):
                parsed = self._artifact_from_canonical_steps(
                    atomized.steps,
                    path,
                    raw_text,
                    source_type,
                    atomized.atomization_mode,
                )
            else:
                parsed = self._atomize_with_rules(raw_text, class_names)
        else:
            parsed = self._atomize_with_rules(raw_text, class_names)

        parsed = self._normalize_artifact(parsed, path, raw_text, source_type)
        parsed.metadata.setdefault(
            "source_refs",
            [ref.to_json() for ref in extraction.refs],
        )
        self.validate_artifact(parsed)
        return parsed

    def export_json(
        self, artifact: SOPDocumentArtifact, destination: str | Path
    ) -> Path:
        dest = Path(destination)
        dest.parent.mkdir(parents=True, exist_ok=True)
        dest.write_text(
            json.dumps(artifact.to_json(), ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return dest

    def _extract_text(self, path: Path) -> str:
        if path.suffix.lower() == ".txt" or path.suffix.lower() == ".md":
            return path.read_text(encoding="utf-8", errors="replace")
        if path.suffix.lower() == ".pptx":
            return self._extract_pptx_text(path)
        return self._extract_pdf_text(path)

    def extract_document(self, source_path: str | Path) -> SOPDocumentExtraction:
        path = Path(source_path)
        if not path.exists():
            raise FileNotFoundError(f"SOP document not found: {path}")
        suffix = path.suffix.lower()
        if suffix not in _DOC_SUFFIXES:
            raise ValueError(f"Unsupported SOP document type: {path.suffix}")

        if suffix in {".txt", ".md"}:
            raw_text = path.read_text(encoding="utf-8", errors="replace")
            refs = self._build_text_refs(raw_text)
        elif suffix == ".pptx":
            refs = self._extract_pptx_refs(path)
            raw_text = "\n\n".join(ref.text for ref in refs if ref.text.strip()).strip()
        else:
            refs = self._extract_pdf_refs(path)
            raw_text = "\n\n".join(ref.text for ref in refs if ref.text.strip()).strip()
            if not raw_text:
                raw_text = self._extract_pdf_text(path)
        title = (
            self._extract_title(
                [ref.label for ref in refs]
                + [line.strip() for line in raw_text.splitlines() if line.strip()]
            )
            or path.stem
        )
        return SOPDocumentExtraction(
            source_path=str(path),
            source_type=suffix.lstrip("."),
            title=title,
            raw_text=raw_text,
            refs=refs,
            metadata={"source_file_name": path.name},
        )

    def _build_text_refs(self, text: str) -> List[SOPSourceRef]:
        refs: List[SOPSourceRef] = []
        blocks = [
            block.strip() for block in re.split(r"\n\s*\n", text) if block.strip()
        ]
        if not blocks:
            blocks = [line.strip() for line in text.splitlines() if line.strip()]
        for idx, block in enumerate(blocks, start=1):
            refs.append(
                SOPSourceRef(
                    kind="section",
                    index=idx,
                    label=f"Section {idx}",
                    text=block,
                )
            )
        return refs

    def _extract_pdf_text(self, path: Path) -> str:
        refs = self._extract_pdf_refs(path)
        text = "\n".join(ref.text for ref in refs if ref.text.strip()).strip()
        if text:
            return text
        return self._ocr_pdf(path)

    def _extract_pdf_refs(self, path: Path) -> List[SOPSourceRef]:
        refs: List[SOPSourceRef] = []
        try:
            from pypdf import PdfReader  # type: ignore[import]

            reader = PdfReader(str(path))
            for index, page in enumerate(reader.pages, start=1):
                try:
                    page_text = (page.extract_text() or "").strip()
                except Exception:
                    page_text = ""
                if page_text:
                    refs.append(
                        SOPSourceRef(
                            kind="page",
                            index=index,
                            label=f"Page {index}",
                            text=page_text,
                        )
                    )
        except Exception:
            refs = []
        if refs:
            return refs
        ocr_text = self._ocr_pdf(path)
        if not ocr_text.strip():
            return []
        return self._build_text_refs(ocr_text)

    def _ocr_pdf(self, path: Path) -> str:
        try:
            import fitz  # type: ignore[import]
        except Exception:
            return ""

        texts: List[str] = []
        doc = fitz.open(str(path))
        try:
            for page in doc:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                img_bytes = pix.tobytes("png")
                text = self._ocr_page_image(img_bytes)
                if text:
                    texts.append(text)
        finally:
            doc.close()
        return "\n".join(texts).strip()

    def _ocr_page_image(self, image_bytes: bytes) -> str:
        if self._ocr is not None:
            try:
                import cv2  # type: ignore[import]
                import numpy as np  # type: ignore[import]

                arr = np.frombuffer(image_bytes, dtype=np.uint8)
                bgr = cv2.imdecode(arr, cv2.IMREAD_COLOR)
                if bgr is not None:
                    scan = getattr(self._ocr, "scan_all", None)
                    if callable(scan):
                        regions = scan(bgr)
                        merged = " ".join(
                            str(getattr(region, "text", "")).strip()
                            for region in regions
                            if str(getattr(region, "text", "")).strip()
                        )
                        if merged.strip():
                            return merged.strip()
            except Exception:
                pass

        try:
            from PIL import Image  # type: ignore[import]
            import pytesseract  # type: ignore[import]

            image = Image.open(io.BytesIO(image_bytes))
            return pytesseract.image_to_string(image).strip()
        except Exception:
            return ""

    def _extract_pptx_text(self, path: Path) -> str:
        refs = self._extract_pptx_refs(path)
        return "\n\n".join(ref.text for ref in refs if ref.text.strip()).strip()

    def _extract_pptx_refs(self, path: Path) -> List[SOPSourceRef]:
        refs: List[SOPSourceRef] = []
        try:
            from pptx import Presentation  # type: ignore[import]
        except Exception:
            return refs

        prs = Presentation(str(path))
        for slide_idx, slide in enumerate(prs.slides, start=1):
            parts: List[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and str(text).strip():
                    parts.append(str(text).strip())
                    continue
                image = getattr(shape, "image", None)
                if image is not None:
                    ocr_text = self._ocr_page_image(image.blob)
                    if ocr_text:
                        parts.append(ocr_text)
            slide_text = "\n".join(part for part in parts if part.strip()).strip()
            if slide_text:
                refs.append(
                    SOPSourceRef(
                        kind="slide",
                        index=slide_idx,
                        label=f"Slide {slide_idx}",
                        text=slide_text,
                    )
                )
        return refs

    def _artifact_from_canonical_steps(
        self,
        canonical_steps: Sequence[Dict[str, Any]],
        path: Path,
        raw_text: str,
        source_type: str,
        atomization_mode: str,
    ) -> Dict[str, Any]:
        steps: List[Dict[str, Any]] = []
        for idx, step in enumerate(canonical_steps, start=1):
            if not isinstance(step, dict):
                continue
            step_type = self._canonical_action_to_type(
                str(step.get("action_kind") or "click")
            )
            target_obj = (
                step.get("target") if isinstance(step.get("target"), dict) else {}
            )
            target_name = (
                str(target_obj.get("name") or target_obj.get("text") or "").strip()
                or None
            )
            atom = SOPAtom(
                id=str(step.get("id") or f"step_{idx:03d}"),
                name=str(step.get("title") or f"Step {idx}"),
                type=step_type,
                target=target_name,
                description=str(step.get("intent") or step.get("title") or ""),
                class_name=target_name,
                confidence=self._safe_float(step.get("confidence")),
                source_span=(
                    str((step.get("source_refs") or [{}])[0].get("label", ""))
                    if step.get("source_refs")
                    else None
                ),
                preconditions=self._ensure_string_list(step.get("preconditions")),
                postconditions=self._ensure_string_list(step.get("postconditions")),
                raw_text=str(step.get("intent") or ""),
            )
            steps.append(atom.to_step())
        title = (
            self._extract_title([step.get("name", "") for step in steps]) or path.stem
        )
        return {
            "version": "6.0.0",
            "title": title,
            "source_path": str(path),
            "source_type": source_type,
            "raw_text": raw_text,
            "metadata": {"atomization_mode": atomization_mode},
            "steps": steps,
        }

    def _canonical_action_to_type(self, action_kind: str) -> str:
        mapping = {
            "click": "click",
            "input": "type_text",
            "wait": "wait_ms",
            "drag": "drag",
            "auth": "auth_sequence",
            "validate": "validate_pins",
        }
        return mapping.get(action_kind, "click")

    def _atomize_with_rules(
        self, raw_text: str, class_names: Sequence[str]
    ) -> Dict[str, Any]:
        steps: List[Dict[str, Any]] = []
        lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
        index = 1
        for line in lines:
            if len(line) < 3:
                continue
            step_type = self._infer_step_type(line)
            class_name = self._infer_class_name(line, class_names)
            step = SOPAtom(
                id=f"step_{index:03d}",
                name=line[:80],
                type=step_type,
                target=class_name,
                description=line,
                class_name=class_name,
                confidence=0.5 if class_name else 0.25,
                raw_text=line,
            )
            steps.append(step.to_step())
            index += 1
        title = self._extract_title(lines) or "Imported SOP"
        return {
            "version": "6.0.0",
            "title": title,
            "source_path": "",
            "source_type": "text",
            "raw_text": raw_text,
            "metadata": {"atomization_mode": "rules"},
            "steps": steps,
        }

    def _normalize_artifact(
        self,
        data: Dict[str, Any],
        path: Path,
        raw_text: str,
        source_type: str,
    ) -> SOPDocumentArtifact:
        steps = self._normalize_steps(data.get("steps", []))
        metadata = data.get("metadata", {})
        if not isinstance(metadata, dict):
            metadata = {"note": str(metadata)}
        metadata.setdefault("source_file_name", path.name)
        metadata.setdefault(
            "atomization_mode", "llm" if self._llm is not None else "rules"
        )
        metadata.setdefault("class_registry", ClassRegistry.load().class_names())
        return SOPDocumentArtifact(
            version=str(data.get("version", "6.0.0")),
            title=str(data.get("title", path.stem)),
            source_path=str(path),
            source_type=source_type,
            raw_text=raw_text,
            steps=steps,
            metadata=metadata,
        )

    def _normalize_steps(self, steps: Any) -> List[Dict[str, Any]]:
        normalized: List[Dict[str, Any]] = []
        if not isinstance(steps, list):
            return normalized
        for idx, item in enumerate(steps, start=1):
            if not isinstance(item, dict):
                continue
            step_type = str(item.get("type", "click")).strip()
            if step_type not in _STEP_TYPES:
                step_type = "click"
            atom = SOPAtom(
                id=str(item.get("id") or f"step_{idx:03d}"),
                name=str(item.get("name") or item.get("description") or f"Step {idx}"),
                type=step_type,
                target=self._sanitize_target(item.get("target")),
                description=str(
                    item.get("description") or item.get("source_text") or ""
                ),
                enabled=bool(item.get("enabled", True)),
                class_name=self._sanitize_target(item.get("class_name")),
                confidence=self._safe_float(item.get("confidence")),
                source_page=self._safe_int(item.get("source_page")),
                source_span=self._safe_text(item.get("source_span")),
                preconditions=self._ensure_string_list(item.get("preconditions")),
                postconditions=self._ensure_string_list(item.get("postconditions")),
                raw_text=str(item.get("source_text") or item.get("raw_text") or ""),
                extra={
                    key: value
                    for key, value in item.items()
                    if key
                    not in {
                        "id",
                        "name",
                        "type",
                        "target",
                        "description",
                        "enabled",
                        "class_name",
                        "confidence",
                        "source_page",
                        "source_span",
                        "preconditions",
                        "postconditions",
                        "source_text",
                        "raw_text",
                    }
                },
            )
            normalized.append(atom.to_step())
        return normalized

    def validate_artifact(self, artifact: SOPDocumentArtifact) -> None:
        if not artifact.steps:
            raise ValueError("SOP document did not produce any steps.")
        for step in artifact.steps:
            if not step.get("id") or not step.get("name"):
                raise ValueError(f"Invalid SOP step: {step!r}")
            if step.get("type") not in _STEP_TYPES:
                raise ValueError(f"Unsupported SOP step type: {step.get('type')!r}")

    def _infer_step_type(self, line: str) -> str:
        lowered = line.lower()
        if any(word in lowered for word in ("wait", "pause", "sleep")):
            return "wait_ms"
        if any(word in lowered for word in ("enter", "type", "input")):
            return "type_text"
        if any(word in lowered for word in ("press", "confirm", "ok", "enter key")):
            return "press_key"
        if any(word in lowered for word in ("drag", "roi", "select area", "box")):
            return "drag"
        if any(word in lowered for word in ("verify", "check", "count")):
            return "validate_pins"
        if any(
            word in lowered
            for word in ("login", "recipe", "save", "apply", "open", "click")
        ):
            return "click"
        return "click"

    def _infer_class_name(self, line: str, class_names: Sequence[str]) -> Optional[str]:
        lowered = line.lower()
        for class_name in class_names:
            if class_name.lower().replace("_", " ") in lowered:
                return class_name
        aliases = {
            "left mold": "mold_left_label",
            "right mold": "mold_right_label",
            "mold left": "mold_left_label",
            "mold right": "mold_right_label",
            "pin": "connector_pin",
            "login": "login_button",
            "recipe": "recipe_button",
            "apply": "apply_button",
            "save": "save_button",
            "open": "open_icon",
            "image source": "image_source",
        }
        for needle, class_name in aliases.items():
            if needle in lowered and class_name in class_names:
                return class_name
        return None

    def _extract_title(self, lines: Sequence[str]) -> str:
        for line in lines[:5]:
            if len(line) > 4:
                return line.strip(" -:\t")
        return ""

    def _sanitize_target(self, value: Any) -> Optional[str]:
        if value is None:
            return None
        text = str(value).strip()
        return text or None

    def _safe_float(self, value: Any) -> Optional[float]:
        try:
            if value is None or value == "":
                return None
            return float(value)
        except Exception:
            return None

    def _safe_int(self, value: Any) -> Optional[int]:
        try:
            if value is None or value == "":
                return None
            return int(value)
        except Exception:
            return None

    def _safe_text(self, value: Any) -> Optional[str]:
        if value is None:
            return None
        text = str(value).strip()
        return text or None

    def _ensure_string_list(self, value: Any) -> List[str]:
        if not value:
            return []
        if isinstance(value, list):
            return [str(item).strip() for item in value if str(item).strip()]
        if isinstance(value, str):
            return [value] if value.strip() else []
        return [str(value)]
